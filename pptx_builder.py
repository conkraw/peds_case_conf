"""PowerPoint export logic for the Pediatric Case Conference Builder."""

from __future__ import annotations

from io import BytesIO
from typing import Any, Dict, Iterable, List

import qrcode
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from case_schema import CASE_SLIDES
from feedback_config import (
    FEEDBACK_DISPLAY_URL,
    FEEDBACK_INSTRUCTION,
    FEEDBACK_QR_URL,
    THANK_YOU_MESSAGE,
    THANK_YOU_TITLE,
)


COLOR_DARK = RGBColor(35, 35, 35)
COLOR_MID = RGBColor(95, 95, 95)
COLOR_LIGHT_GRAY = RGBColor(242, 242, 242)
COLOR_HEADER = RGBColor(70, 70, 70)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_ACCENT = RGBColor(30, 80, 130)
COLOR_ACCENT_LIGHT = RGBColor(226, 236, 246)
COLOR_PROMPT_LIGHT = RGBColor(250, 238, 218)
COLOR_GREEN_LIGHT = RGBColor(226, 242, 228)

SLIDE_W = 13.333
SLIDE_H = 7.5


def _safe_text(value: Any) -> str:
    return "" if value is None else str(value).strip()


def _lines(value: Any) -> List[str]:
    return [line.strip() for line in _safe_text(value).splitlines() if line.strip()]


def _first_line(value: Any) -> str:
    lines = _lines(value)
    return lines[0] if lines else ""


def add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: Any,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = COLOR_DARK,
    align=PP_ALIGN.LEFT,
    fill: RGBColor | None = None,
    margin: float = 0.08,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = _safe_text(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color

    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = fill

    return shape


def add_title(slide, title: str, subtitle: str | None = None):
    add_textbox(slide, 0.55, 0.20, 12.2, 0.52, title, font_size=29, bold=True, color=COLOR_DARK)
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(0.82), Inches(12.1), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.color.rgb = COLOR_ACCENT
    if subtitle:
        add_textbox(slide, 0.65, 0.90, 12.0, 0.35, subtitle, font_size=14, color=COLOR_MID)


def add_footer(slide, text: str = "Case Conference Builder"):
    add_textbox(slide, 0.6, 7.10, 12.1, 0.22, text, font_size=8, color=COLOR_MID, align=PP_ALIGN.RIGHT)


def add_bullets(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    items: Iterable[str],
    font_size: int = 17,
    bullet: bool = True,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    usable_items = [str(item).strip() for item in items if str(item).strip()]
    if not usable_items:
        usable_items = [""]
    for i, item in enumerate(usable_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix_ok = item.startswith(("•", "-", "A.", "B.", "C.", "D.", "E.", "1.", "2.", "3."))
        p.text = f"• {item}" if bullet and not prefix_ok else item
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_DARK
        p.level = 0
    return shape


def add_section_label(slide, x: float, y: float, w: float, label: str, fill: RGBColor = COLOR_ACCENT_LIGHT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.bold = True
    r.font.size = Pt(14)
    r.font.color.rgb = COLOR_DARK
    return shape


def add_content_block(slide, x: float, y: float, w: float, h: float, label: str, value: Any, *, prompt: bool = False):
    fill = COLOR_PROMPT_LIGHT if prompt else COLOR_ACCENT_LIGHT
    add_section_label(slide, x, y, w, label, fill=fill)
    text_y = y + 0.40
    text_h = max(0.25, h - 0.42)
    lines = _lines(value)
    if len(lines) > 1:
        return add_bullets(slide, x, text_y, w, text_h, lines, font_size=15 if len(lines) > 5 else 16)
    return add_textbox(slide, x, text_y, w, text_h, _safe_text(value), font_size=16, color=COLOR_DARK)


def add_table(slide, rows_data: List[Dict[str, Any]], columns: List[str], x=0.55, y=1.45, w=12.25, h=4.05):
    cleaned_rows = []
    for row in rows_data or []:
        if any(_safe_text(row.get(col, "")) for col in columns):
            cleaned_rows.append({col: _safe_text(row.get(col, "")) for col in columns})
    cleaned_rows = cleaned_rows[:6]
    if not cleaned_rows:
        cleaned_rows = [{col: "" for col in columns}]

    shape = slide.shapes.add_table(
        len(cleaned_rows) + 1,
        len(columns),
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    table = shape.table
    col_width = w / max(len(columns), 1)
    for idx in range(len(columns)):
        table.columns[idx].width = Inches(col_width)

    for c, col in enumerate(columns):
        cell = table.cell(0, c)
        cell.text = col
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_HEADER
        cell.margin_left = Inches(0.04)
        cell.margin_right = Inches(0.04)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(10)
                run.font.color.rgb = COLOR_WHITE

    for r, row in enumerate(cleaned_rows, start=1):
        for c, col in enumerate(columns):
            cell = table.cell(r, c)
            cell.text = _safe_text(row.get(col, ""))
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_LIGHT_GRAY
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(8.5 if len(columns) >= 4 else 9.5)
                    run.font.color.rgb = COLOR_DARK
    return shape


def _slide_subtitle(deck: Dict[str, Dict[str, Any]]) -> str:
    title_data = deck.get("title_goal", {})
    case_title = _safe_text(title_data.get("case_title"))
    presenter = _safe_text(title_data.get("presenter_name"))
    if case_title and presenter:
        return f"{case_title} | {presenter}"
    return case_title or presenter


def _is_prompt_slide(slide_id: str) -> bool:
    return slide_id in {"history_questions", "exam_questions", "diagnostic_studies_question"}


def _build_title_slide(prs: Presentation, deck: Dict[str, Dict[str, Any]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    data = deck.get("title_goal", {})
    case_title = _safe_text(data.get("case_title")) or "Case Conference"
    presenter = _safe_text(data.get("presenter_name"))
    learning_point = _safe_text(data.get("learning_point"))
    session_goal = _safe_text(data.get("session_goal"))
    reasoning_skill = _safe_text(data.get("reasoning_skill"))
    setting = _safe_text(data.get("case_setting"))

    add_textbox(slide, 0.7, 0.55, 11.9, 0.75, "Pediatric Case Conference", font_size=34, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.1, 1.45, 11.1, 0.70, case_title, font_size=28, bold=True, color=COLOR_DARK, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.2, 2.20, 10.9, 0.35, f"Presenter: {presenter}" if presenter else "", font_size=16, color=COLOR_MID, align=PP_ALIGN.CENTER)
    if setting:
        add_textbox(slide, 1.2, 2.55, 10.9, 0.35, setting, font_size=14, color=COLOR_MID, align=PP_ALIGN.CENTER)

    add_content_block(slide, 1.15, 3.10, 11.05, 1.00, "Learning Point", learning_point)
    add_content_block(slide, 1.15, 4.30, 11.05, 1.10, "Session Goal", session_goal)
    add_textbox(slide, 1.15, 5.78, 5.4, 0.42, f"Reasoning focus: {reasoning_skill}", font_size=14, bold=True, color=COLOR_ACCENT, fill=COLOR_ACCENT_LIGHT, align=PP_ALIGN.CENTER)
    add_textbox(slide, 6.80, 5.78, 5.4, 0.42, "Progressive case reveal | Audience reasoning | Bedside meaning", font_size=13, bold=True, color=COLOR_DARK, fill=COLOR_GREEN_LIGHT, align=PP_ALIGN.CENTER)
    add_footer(slide)


def _build_generic_slide(prs: Presentation, slide_def: Dict[str, Any], deck: Dict[str, Dict[str, Any]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_id = slide_def["id"]
    slide_data = deck.get(slide_id, {})
    add_title(slide, _safe_text(slide_def.get("export_title") or slide_def.get("label")), subtitle=_slide_subtitle(deck))

    public_fields = [f for f in slide_def.get("fields", []) if not f.get("private_note")]
    table_fields = [f for f in public_fields if f.get("type") == "table"]
    text_fields = [f for f in public_fields if f.get("type") != "table"]
    y = 1.35

    if table_fields:
        table_field = table_fields[0]
        # Put a short text block above the table when present.
        pre_text = text_fields[:1]
        if pre_text:
            add_content_block(
                slide,
                0.60,
                1.22,
                12.10,
                0.95,
                pre_text[0]["label"],
                slide_data.get(pre_text[0]["key"], ""),
                prompt=_is_prompt_slide(slide_id),
            )
            table_y = 2.35
            table_h = 3.35
        else:
            table_y = 1.45
            table_h = 4.15
        add_table(slide, slide_data.get(table_field["key"], []), table_field.get("columns", []), y=table_y, h=table_h)
        remaining = text_fields[1:] if pre_text else text_fields
        if remaining:
            block_count = min(len(remaining), 2)
            block_w = 12.1 / block_count
            for idx, field in enumerate(remaining[:2]):
                add_content_block(
                    slide,
                    0.60 + idx * block_w,
                    5.95,
                    block_w - 0.12,
                    0.95,
                    field["label"],
                    slide_data.get(field["key"], ""),
                )
    else:
        count = max(len(text_fields), 1)
        # If a slide has several fields, use two columns. Otherwise use roomy vertical blocks.
        if count >= 4:
            positions = [
                (0.65, 1.30, 5.95, 1.25),
                (6.75, 1.30, 5.95, 1.25),
                (0.65, 2.85, 5.95, 1.65),
                (6.75, 2.85, 5.95, 1.65),
                (0.65, 4.85, 5.95, 1.75),
                (6.75, 4.85, 5.95, 1.75),
            ]
        elif count == 3:
            positions = [(0.70, 1.35, 11.95, 1.35), (0.70, 2.95, 11.95, 1.65), (0.70, 4.85, 11.95, 1.70)]
        elif count == 2:
            positions = [(0.75, 1.55, 11.85, 2.15), (0.75, 4.05, 11.85, 2.35)]
        else:
            positions = [(0.85, 1.70, 11.65, 4.70)]

        for field, pos in zip(text_fields[: len(positions)], positions):
            add_content_block(
                slide,
                *pos,
                field["label"],
                slide_data.get(field["key"], ""),
                prompt=_is_prompt_slide(slide_id),
            )
    add_footer(slide)


def _build_feedback_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, THANK_YOU_TITLE)
    add_textbox(slide, 0.85, 1.45, 7.1, 0.8, THANK_YOU_MESSAGE, font_size=21, color=COLOR_DARK)
    add_textbox(slide, 0.85, 2.45, 7.1, 0.6, FEEDBACK_INSTRUCTION, font_size=18, bold=True, color=COLOR_ACCENT)
    add_textbox(slide, 0.85, 3.25, 7.1, 0.6, FEEDBACK_DISPLAY_URL, font_size=18, color=COLOR_DARK, fill=COLOR_ACCENT_LIGHT)

    qr_img = BytesIO()
    qrcode.make(FEEDBACK_QR_URL).save(qr_img, format="PNG")
    qr_img.seek(0)
    slide.shapes.add_picture(qr_img, Inches(9.1), Inches(1.65), width=Inches(2.8), height=Inches(2.8))
    add_textbox(slide, 8.45, 4.70, 4.1, 0.35, "Feedback QR code", font_size=13, color=COLOR_MID, align=PP_ALIGN.CENTER)
    add_footer(slide)


def _build_facilitator_notes_slide(prs: Presentation, deck: Dict[str, Dict[str, Any]]):
    notes: List[str] = []
    for slide_def in CASE_SLIDES:
        slide_data = deck.get(slide_def["id"], {})
        for field in slide_def.get("fields", []):
            if field.get("private_note"):
                text = _safe_text(slide_data.get(field["key"], ""))
                if text:
                    notes.append(f"{slide_def.get('export_title', slide_def['label'])}: {text}")
    if not notes:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Facilitator Notes")
    add_bullets(slide, 0.75, 1.35, 11.85, 5.65, notes[:8], font_size=12, bullet=True)
    add_footer(slide)


def build_powerpoint(deck: Dict[str, Dict[str, Any]], include_facilitator_notes: bool = True) -> BytesIO:
    """Build the case conference PowerPoint and return it as a BytesIO object."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    _build_title_slide(prs, deck)
    for slide_def in CASE_SLIDES[1:]:
        _build_generic_slide(prs, slide_def, deck)
    if include_facilitator_notes:
        _build_facilitator_notes_slide(prs, deck)
    _build_feedback_slide(prs)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output
